
from pptx import Presentation
from pptx.dml.color import RGBColor
from PIL import Image
from pptx.util import Inches 
from pptx.util import Pt
import os
import numpy as np
import json


with open('settings.json') as  f:
    settings = json.load(f)
PRS_WIDTH_INCHES = settings['presentation_width_inches']
PRS_HEIGHT_INCHES = settings['presentation_height_inches']
PIXELS_PER_INCH = settings['pixels_per_inch']

def text_in_presentation(presntation,nr_slide , nr_shape , text,font_size= 66,font_name = 'Calibri'):
    text_frame = presntation.slides[nr_slide].shapes[nr_shape].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    font = run.font
    font.size = Pt(font_size)
    font.name = font_name
    font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.text = text

def requaried_new_size_image(path_of_image , max_height = (PRS_HEIGHT_INCHES -1) * PIXELS_PER_INCH, max_width= (PRS_WIDTH_INCHES - 1 ) * PIXELS_PER_INCH):
    """Function give optimal size wich can be on presentation with keeping size ratio
    @pam path_of_image: str path to image
    @pam max_width: int max width of image in pixels
    @pam max_height: int max height of image in pixels
    @return: new_height , new_width: tuple with optimal size in Inches"""
    print('Wwołanie')

    height, width, _ = np.array(Image.open(path_of_image)).shape

    if height >= max_height and width >= max_width:
        print('Tutaj 1')
        if height/ max_height < width / max_width:
            print('Wersja 1')
            new_height = max_height
            new_width = max_width * height/ width
        else:
            print('Wersja 2')
            new_width = max_width
            new_height = max_height * width/ height
    elif height >= max_height:
        print('Tutaj 2')
        print(height)
        new_height = max_height
        new_width = max_width * height/ width
        print(new_height)
    elif width >= max_width:
        print('Tutaj 3')
        new_width = max_width
        new_height = max_height * width/ height
    else:
        print('Tutaj 4')
        if height/ max_height < width / max_width:
            print('Przypadek 1')
            new_height = max_height
            new_width = max_width * height/ width
        else:
            print('Przypadek 2')
            new_width = max_width
            new_height = max_height * width/ height
    return new_height / PIXELS_PER_INCH , new_width / PIXELS_PER_INCH

def make_presentation():
    prs = Presentation('templates\Templatka 1 (Orla na Piątkę).pptx')


    level = 6
    name_of_route = ' Dolina 5 stawów'
    heightest_point = 'Chełmiec, 851 m n.p.m.'
    elevation = '600 m'
    up = '700 m'
    down = '701 m'
    distance = '20 km'
    time = '7h'
    leader = 'Maciej Karczewski'
    path_profile = 'zdjecie.png'
    path_route = 'route.jpg'




    # text_in_presentation(prs,nr_slide =0,nr_shape=0,text=f'{name_of_route} \n Poziom trudności: {level} ',font_size=66)
    # text_in_presentation(prs,nr_slide =1,nr_shape=3,text=f'Najwyższy szczyt: {heightest_point} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=4,text=f'Przewyższenie: {elevation} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=5,text=f'Podejście: {up} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=6,text=f'Zejście: {down} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=7,text=f'Szacowany czas przejścia: {time} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=8,text=f'Długość: {distance} ',font_size=28)
    # text_in_presentation(prs,nr_slide =6,nr_shape=0,text=f'Prowadzi \n {leader}',font_size=66)


    text_in_presentation(prs,nr_slide =0,nr_shape=0,text=f'{name_of_route} \n Poziom trudności: {level} ',font_size=66)
    text_in_presentation(prs,nr_slide =1,nr_shape=2,text=f'Najwyższy szczyt: {heightest_point} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=3,text=f'Przewyższenie: {elevation} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=4,text=f'Podejście: {up} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=5,text=f'Zejście: {down} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=6,text=f'Szacowany czas przejścia: {time} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=7,text=f'Długość: {distance} ',font_size=28)
    text_in_presentation(prs,nr_slide =6,nr_shape=0,text=f'Prowadzi \n {leader}',font_size=66)


    # Add profile to presentation
    prs.slides[3].shapes.add_picture(path_profile,left=Inches(1),top=Inches(1), width=Inches(PRS_WIDTH_INCHES - 2), height=Inches(5))


    # First image of views 
    height , width = requaried_new_size_image(settings['first_image'],max_width= (PRS_WIDTH_INCHES -1 )* PIXELS_PER_INCH / 2)
    left = (PRS_WIDTH_INCHES/ 2 - width) /2 + PRS_WIDTH_INCHES/2
    top = (PRS_HEIGHT_INCHES - height )/ 2
    prs.slides[1].shapes.add_picture(settings['first_image'], left= Inches(left), top= Inches(top), width = Inches(width), height = Inches(height))

    # Second and third images of views
    height , width = requaried_new_size_image(settings['second_image'],max_height = (PRS_HEIGHT_INCHES -2.5) * PIXELS_PER_INCH, max_width= (PRS_WIDTH_INCHES -1 )* PIXELS_PER_INCH / 2)
    left = (PRS_WIDTH_INCHES/ 2 - width) /2 
    top = (PRS_HEIGHT_INCHES + 1.5  - height )/ 2 
    prs.slides[2].shapes.add_picture(settings['second_image'], left= Inches(left), top= Inches(top), width = Inches(width), height = Inches(height))

    height , width = requaried_new_size_image(settings['third_image'],max_height = (PRS_HEIGHT_INCHES -2.5) * PIXELS_PER_INCH,max_width= (PRS_WIDTH_INCHES -1 )* PIXELS_PER_INCH / 2)
    left = (PRS_WIDTH_INCHES/ 2 - width) /2 + PRS_WIDTH_INCHES/2
    top = (PRS_HEIGHT_INCHES + 1.5 - height )/ 2 
    prs.slides[2].shapes.add_picture(settings['third_image'], left= Inches(left), top= Inches(top), width = Inches(width), height = Inches(height))


    # Image of route
    height , width = requaried_new_size_image(path_route)
    left = (PRS_WIDTH_INCHES - width )/ 2
    top = (PRS_HEIGHT_INCHES - height )/ 2
    prs.slides[4].shapes.add_picture(path_route, left=Inches(left),top=Inches(top), width= Inches(width) , height=Inches(height))
    # prs.slides[4].shapes.add_picture(path_route, left=Inches(1),top=Inches(1))
    # width=Inches(11), height=Inches(7)
    # prs.slides[4].shapes.add_picture(path_route, left=Inches(1),top=Inches(1), width= Inches(12) )


    prs.save("testy.pptx")
    os.system("testy.pptx")
if __name__ == "__main__":
    make_presentation()







