
from importlib.resources import path
from tkinter import W
from pptx import Presentation
from pptx.dml.color import RGBColor
from PIL import Image
from pptx.util import Inches 
from pptx.util import Pt
import os
import numpy as np
import json


# with open(os.path.join('settings.json')) as  f:
#     settings = json.load(f)
# PRS_WIDTH_INCHES = settings['presentation_width_inches']
# PRS_HEIGHT_INCHES = settings['presentation_height_inches']
# PIXELS_PER_INCH = settings['pixels_per_inch']
PRS_WIDTH_INCHES = 13.33
PRS_HEIGHT_INCHES = 7.5
PIXELS_PER_INCH = 72

def text_in_presentation(presentation,nr_slide , nr_shape , text,font_size= 66,font_name = 'Calibri'):
    text_frame = presentation.slides[nr_slide].shapes[nr_shape].text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    font = run.font
    font.size = Pt(font_size)
    font.name = font_name
    font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.text = text

def required_new_size_image(path_of_image , max_height = (PRS_HEIGHT_INCHES -1) * PIXELS_PER_INCH, max_width= (PRS_WIDTH_INCHES - 1 ) * PIXELS_PER_INCH):
    """Function give optimal size which can be on presentation with keeping size ratio
    @pam path_of_image: str path to image
    @pam max_width: int max width of image in pixels
    @pam max_height: int max height of image in pixels
    @return: new_height , new_width: tuple with optimal size in Inches"""
    print(max_height, max_width)

    height, width, _ = np.array(Image.open(path_of_image)).shape

    if height >= max_height and width >= max_width:
        print('Tutaj 1')
        if height/ max_height < width / max_width:
            print('Wersja 1')
            new_height = max_height
            new_width = width *  max_height / height
        else:
            print('Wersja 2')
            new_width = max_width
            new_height = height *  max_width / width
    elif height >= max_height:
        print('Tutaj 2')
        print(height)
        new_height = max_height
        new_width = width *  max_height / height
        print(new_height)
    elif width >= max_width:
        print('Tutaj 3')
        new_width = max_width
        new_height = height * max_width/ width
    else:
        print('Tutaj 4')
        if height/ max_height < width / max_width:
            print('Przypadek 1')
            new_height = max_height
            new_width = width * max_height/ height
        else:
            print('Przypadek 2')
            new_width = max_width
            new_height = height * max_width / width
    return new_height / PIXELS_PER_INCH , new_width / PIXELS_PER_INCH

def make_presentation():

    
    
    with open('data.json','r',encoding="utf-8") as f:
        data = json.load(f)
    with open(os.path.join('../../','settings.json'),'r',encoding='utf-8') as f:
        settings = json.load(f)
    print(os.getcwd())
    # print([name for name in os.listdir(os.path.join("../../","templates")) if data["background"]  in name and ".pptx" in name])
    # print([name for name in os.listdir(os.path.join("..","templates")) if data["background"]  in name and ".pptx" in name][0])
    path_of_template = os.path.join('../../','templates', f'{[name for name in os.listdir(os.path.join("../..","templates")) if data["background"]  in name and ".pptx" in name][0]}')
    print(path_of_template)
    prs = Presentation(path_of_template)

    path_profile = settings["path_profile"]
    path_route = settings["path_route"]




    # text_in_presentation(prs,nr_slide =0,nr_shape=0,text=f'{name_of_route} \n Poziom trudności: {level} ',font_size=66)
    # text_in_presentation(prs,nr_slide =1,nr_shape=3,text=f'Najwyższy szczyt: {heightest_point} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=4,text=f'Przewyższenie: {elevation} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=5,text=f'Podejście: {up} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=6,text=f'Zejście: {down} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=7,text=f'Szacowany czas przejścia: {time} ',font_size=28)
    # text_in_presentation(prs,nr_slide =1,nr_shape=8,text=f'Długość: {distance} ',font_size=28)
    # text_in_presentation(prs,nr_slide =6,nr_shape=0,text=f'Prowadzi \n {leader}',font_size=66)
    # text_in_presentation(prs,nr_slide =5,nr_shape=0,text=f'Uwagi:',font_size=56)


    text_in_presentation(prs,nr_slide =0,nr_shape=0,text=f'{data["name_of_route"]} \n Poziom trudności: {data["level"]} ',font_size=66)
    text_in_presentation(prs,nr_slide =1,nr_shape=2,text=f'Najwyższy szczyt: {data["heightest_point"]} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=3,text=f'Przewyższenie: {data["elevation"]} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=4,text=f'Podejście: {data["up"]} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=5,text=f'Zejście: {data["down"]} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=6,text=f'Szacowany czas przejścia: {data["time"]} ',font_size=28)
    text_in_presentation(prs,nr_slide =1,nr_shape=7,text=f'Długość: {data["distance"]} ',font_size=28)
    text_in_presentation(prs,nr_slide =6,nr_shape=0,text=f'Prowadzi \n {data["leader"]}',font_size=66)
    text_in_presentation(prs,nr_slide =5,nr_shape=0,text=f'Uwagi:',font_size=56)



    # add cautions
    prs.slides[5].shapes.add_textbox(left = Inches(1.5), top=Inches(2.2), width = Inches(10), height = Inches(4.5))
    text_frame = prs.slides[5].shapes[1].text_frame
    text_frame.clear()
    for  text in data["coutions"]:
        p = text_frame.add_paragraph()
        run = p.add_run()
        font = run.font
        font.size = Pt(48)
        font.name = 'Calibra'
        font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.text = f'- {text} '


    # Add profile to presentation
    prs.slides[3].shapes.add_picture(path_profile,left=Inches(1),top=Inches(1), width=Inches(PRS_WIDTH_INCHES - 2), height=Inches(5))


    # First image of views 
    try:
        height , width = required_new_size_image(settings['first_view'],max_width= (PRS_WIDTH_INCHES -1 )* PIXELS_PER_INCH / 2)
        left = (PRS_WIDTH_INCHES/ 2 - width) /2 + PRS_WIDTH_INCHES/2
        top = (PRS_HEIGHT_INCHES - height )/ 2
        prs.slides[1].shapes.add_picture(settings['first_view'], left= Inches(left), top= Inches(top), width = Inches(width), height = Inches(height))
    except:
        pass
    # Second and third images of views
    try:
        height , width = required_new_size_image(settings['second_view'],max_height = (PRS_HEIGHT_INCHES -2.5) * PIXELS_PER_INCH, max_width= (PRS_WIDTH_INCHES -2 )* PIXELS_PER_INCH / 2)
        left = (PRS_WIDTH_INCHES/ 2 - width) /2 
        top = (PRS_HEIGHT_INCHES + 1.5  - height )/ 2 
        prs.slides[2].shapes.add_picture(settings['second_view'], left= Inches(left), top= Inches(top), width = Inches(width), height = Inches(height))
    except:
        pass

    try:
        height , width = required_new_size_image(settings['third_view'],max_height = (PRS_HEIGHT_INCHES -2.5) * PIXELS_PER_INCH,max_width= (PRS_WIDTH_INCHES -2 )* PIXELS_PER_INCH / 2)
        left = (PRS_WIDTH_INCHES/ 2 - width) /2 + PRS_WIDTH_INCHES/2
        top = (PRS_HEIGHT_INCHES + 1.5 - height )/ 2 
        prs.slides[2].shapes.add_picture(settings['third_view'], left= Inches(left), top= Inches(top), width = Inches(width), height = Inches(height))
    except:
        pass 

    # Image of route
    height , width = required_new_size_image(path_route)
    left = (PRS_WIDTH_INCHES - width )/ 2
    top = (PRS_HEIGHT_INCHES - height )/ 2
    prs.slides[4].shapes.add_picture(path_route, left=Inches(left),top=Inches(top), width= Inches(width) , height=Inches(height))
    # prs.slides[4].shapes.add_picture(path_route, left=Inches(1),top=Inches(1))
    # width=Inches(11), height=Inches(7)
    # prs.slides[4].shapes.add_picture(path_route, left=Inches(1),top=Inches(1), width= Inches(12) )


    prs.save(f'{data["file_name"]}.pptx')
    npath = data["file_name"].replace(' ','\\ ')
    os.system(f'{npath}.pptx')
if __name__ == "__main__":
    make_presentation()







